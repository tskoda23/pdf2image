import os
import sys
from pgmagick import Image
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfFileWriter, PdfFileReader
import logging

class PdfToPpt(object):
    def __init__(self, pdf_file=None, ppt_file=None):
        self.pdf_file = pdf_file
        self.ppt_file = pdf_file.replace('.pdf', '.pptx')
        self.total_pages = 1    
        self.log = logging.getLogger('PdfToPptx')
        self.log.debug('{} \n {}'.format(self.pdf_file, self.ppt_file))

    def check_file_exist(self, file_path):
        self.log.info('Checking file - %s ' % file_path)
        if os.path.isfile(file_path):
            return True
        else:
            return False

    def pdf_to_image(self, pdf_file):
        if not self.check_file_exist(pdf_file):
            self.log.debug('Requested file not found in {} '.format(pdf_file))
            return False
        image_file = pdf_file.replace('.pdf', '.jpg')
        try:
            pdf_to_img = Image()
            pdf_to_img.density('200')
            pdf_to_img.read(pdf_file)
            pdf_to_img.write(image_file)
            self.log.info('Image convert passed - {} '.format(image_file))
            return True
        except Exception:
            self.log.error('Image convert failed - {} '.format(image_file))
            self.log.error('', exc_info=True)
            return False

    def pdf_splitter(self):
        self.log.info('Called pdf_splitter')
        input_pdf = PdfFileReader(open(self.pdf_file, 'rb'), strict=False)
        self.total_pages = input_pdf.numPages

        for page_number in range(self.total_pages):
            output = PdfFileWriter()
            output.addPage(input_pdf.getPage(page_number))
            # new filename
            new_pdf = '_%s%s' % (str(page_number+1), '.pdf')
            new_pdf = self.pdf_file.replace('.pdf', new_pdf)
            file_stream = open(new_pdf, 'wb')
            output.write(file_stream)
            file_stream.close()

            # calling pdf to image conversion
            self.pdf_to_image(new_pdf)

    def create_ppt(self):
        self.log.info('Called create_ppt')
        prs = Presentation()
        try:
            for slide_number in range(self.total_pages):
                img_path = self.pdf_file.replace('.pdf', '_%s%s' % (str(slide_number+1), '.jpg'))
                self.log.debug('%s' % img_path)
                new_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(new_slide)
                subtitle = slide.placeholders[1]
                title = slide.shapes.title
                title.text = "Image %s " % str(slide_number+1)
                left = top = Inches(0.1)
                height = Inches(7.5)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                prs.save(self.ppt_file)
        except IOError:
            self.log.error('error creating ppt', exc_info=True)

    def execute(self):
        self.log.info('Calling the main execution for ppt conversion')
        self.pdf_splitter()
        self.create_ppt()
        self.log.info('Done ppt conversion')
  


if __name__ == '__main__':
    directory = sys.argv[1]
    quality = int(sys.argv[2])
    file_list = (f for f in os.listdir(directory) if f.endswith('.' + 'pdf'))
    output_directory = '{}\{}'.format(directory, 'OUT') 
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)
    
    for file in file_list:
        full_name = '{}\{}'.format(directory, file)
        PdfToPpt(pdf_file=full_name).execute()


    image_list = (f for f in os.listdir(directory) if f.endswith('.' + 'jpg'))
    for image in image_list:
        full_name = '{}\{}'.format(directory, image)

        real_image = Image(full_name)
        real_image.quality(quality)
        real_image.write(full_name)
        command = 'img2pdf -o {}\{} {}'.format(output_directory, image.replace("_1.jpg", ".pdf"), full_name)
        os.system(command)
